from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .docx import StoreImage


def pptx_to_markdown(data: bytes, store_image: StoreImage) -> str:
    try:
        presentation = Presentation(BytesIO(data))
    except Exception as error:
        raise ValueError(f"cannot open pptx: {error}") from error
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        heading = f"## Slide {index}" + (f" — {title}" if title else "")
        lines = [heading]
        for shape in slide.shapes:
            if slide.shapes.title is not None and shape.shape_id == slide.shapes.title.shape_id:
                continue
            if shape.has_text_frame:
                for line in shape.text_frame.text.splitlines():
                    text = line.strip()
                    if text:
                        lines.append(f"- {text}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    picture = shape.image
                except (AttributeError, ValueError):
                    picture = None
                if picture is not None:
                    image_id = store_image(picture.blob, picture.content_type)
                    lines.append(f"![slide image](ca-image://{image_id})")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("")
                for line in notes.splitlines():
                    if line.strip():
                        lines.append(f"> {line.strip()}")
        sections.append("\n".join(lines))
    if not sections:
        raise ValueError("pptx contains no slides")
    return "\n\n".join(sections) + "\n"
